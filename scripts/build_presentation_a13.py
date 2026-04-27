# -*- coding: utf-8 -*-
"""
Дубликат презентации: Calibri, фото public/projects/pptx-*.webp (или .png), тексты с сайта.
Исходник: c:\\Users\\vladi\\Downloads\\Untitled.pptx
Результат: c:\\Users\\vladi\\Downloads\\Бюро-А13-презентация-Calibri.pptx
"""
from __future__ import annotations

import shutil
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
IMG_DIR = ROOT / "public" / "projects"
SRC_PPTX = Path(r"c:\Users\vladi\Downloads\Untitled.pptx")
OUT_PPTX = Path(r"c:\Users\vladi\Downloads\Бюро-А13-презентация-Calibri.pptx")

PHONE = "+7 (916) 117-13-50"
EMAIL = "info@a13bureau.ru"
ADDRESS_OFFICE = "г. Москва, Рублевское шоссе д.26 корп.4"
ADDRESS_PROD = "г. Фрязино, ул. Горького д.10 стр.1"

ABOUT_BODY = (
    "Бюро А13 — инжиниринговая компания полного цикла с опытом работы в области "
    "конструирования, производства и монтажа светопрозрачных конструкций любой сложности.\n\n"
    "Команда с 2009 года: 226+ реализованных проектов и полный цикл от проекта до монтажа. "
    "Профильные инженеры и монтаж — от коттеджей до инфраструктуры и премиальной застройки.\n\n"
    "Специализируемся на интересных и значимых проектах в любом регионе. "
    "Собственное производство во Фрязино (до 8 000 м² фасадов в месяц) и квалифицированный персонал."
)

SPECS_LINES = [
    ("Светопрозрачные фасады", "Стоечно-ригельные, структурные и полуструктурные решения; расчёты и КМД."),
    ("Алюминиевые окна и двери", "Тёплый и холодный контур; панорамное остекление и входные системы."),
    ("Зенитные фонари и атриумы", "Атриумы и кровельные фонари; ТЦ, офисы и инфраструктура."),
    ("Входные группы", "Алюминиевые двери, козырьки, тамбуры, витрины; полный цикл работ."),
    ("Вентилируемые фасады", "НВФ на подсистеме: АКП, керамогранит, натуральный камень."),
    ("Проектирование", "ОПР, КМ, КМД, ППР; статика и теплотехника, аэродинамика, надзор."),
]

# img_seed: как на сайте — смещение для pptxAsset(projectIndex * 3)
PROJECT_CARDS = [
    {
        "year": "2024",
        "title": "Объект «4 ветра», Белорусская",
        "body": "Поставка и монтаж светопрозрачных конструкций в Москве: фасадные системы и панорамное остекление.",
        "spec": (
            "● Направление: фасады и СПК\n"
            "● Полный цикл: проектирование, производство, монтаж\n"
            "● Системы: по расчёту и ТЗ заказчика (Schüco, Alutech, Reynaers и др.)\n"
            "● Контроль: КМД, заводская сборка, сопровождение на площадке\n"
            "● Объём и сроки: по запросу менеджеров"
        ),
        "img_seed": 6 * 3,
    },
    {
        "year": "2023",
        "title": "ЖК «Зиларт»",
        "body": "Остекление корпусов жилого квартала: фасадные и оконные системы, витражи.",
        "spec": (
            "● Сегмент: жилая недвижимость\n"
            "● Решения: окна, витражи, фасадные узлы\n"
            "● Серийные корпуса: логистика и монтаж по графику\n"
            "● Спецификация: по проекту заказчика\n"
            "● Уточнение параметров: Бюро А13"
        ),
        "img_seed": 17 * 3,
    },
    {
        "year": "2024",
        "title": "Электродепо метро (Саларьево, Руднёво, Выхино, Лихоборы)",
        "body": "Инфраструктура Московского метрополитена; заказчик Мосинжпроект. Зенитные фонари и СПК.",
        "spec": (
            "● Направление: зенитные фонари / инфраструктура\n"
            "● Регламенты: дисциплина на площадке, согласования\n"
            "● Надёжность: поставка, монтаж, сдача участков\n"
            "● Инженерия: расчёты и узлы под объект\n"
            "● Референс: объекты Мосинжпроект"
        ),
        "img_seed": 14 * 3,
    },
    {
        "year": "2022",
        "title": "«Легенда Цветного», Capital Group",
        "body": "Остекление первых этажей и торговых витрин; входные зоны на городском фасаде.",
        "spec": (
            "● Направление: входные группы и витрины\n"
            "● Заказчик: Capital Group\n"
            "● Фокус: качество витражей и входных зон\n"
            "● Городской контекст: плотная застройка\n"
            "● Полный цикл работ Бюро А13"
        ),
        "img_seed": 15 * 3,
    },
    {
        "year": "2021",
        "title": "Реконструкция фасадов ЩЛЗ",
        "body": "Реконструкция фасадов промышленного объекта; заказчик Политехстрой. Вентфасад и СПК.",
        "spec": (
            "● Направление: вентилируемые фасады\n"
            "● Промышленный объект: усиленный контроль узлов\n"
            "● График: в соответствии с договором\n"
            "● Коммуникация: инженерный уровень\n"
            "● Утепление и облицовка по проекту"
        ),
        "img_seed": 16 * 3,
    },
    {
        "year": "2024",
        "title": "Жилой дом, Пионерская 30к10",
        "body": "Остекление для РКК «Энергия»: панорамные решения и безопасные системы.",
        "spec": (
            "● Сегмент: жилой дом / панорамное остекление\n"
            "● Решения: окна и раздвижные системы по расчёту\n"
            "● Теплотехника и акустика: подбор стеклопакетов\n"
            "● Монтаж: собственные бригады\n"
            "● Детали: уточняйте у менеджеров"
        ),
        "img_seed": 10 * 3,
    },
]

REVIEWS = [
    {
        "quote": "Слаженная работа проектного блока и монтажа на кампусе. Сроки и качество остекления соответствовали договору.",
        "initials": "ФС",
        "name": "Представитель заказчика",
        "company": "Фонд «Сколково»",
    },
    {
        "quote": "Профессиональный подход к узлам примыкания и документации. Рекомендуем для объектов с повышенными требованиями.",
        "initials": "РП",
        "name": "Руководитель проекта",
        "company": "ПАО «Лукойл»",
    },
    {
        "quote": "Участие в остеклении инфраструктуры метро: дисциплина на площадке, понимание регламентов.",
        "initials": "ГИ",
        "name": "Главный инженер",
        "company": "Мосинжпроект",
    },
    {
        "quote": "Качественное исполнение витрин и входных зон на сложном городском фасаде.",
        "initials": "ДС",
        "name": "Директор по строительству",
        "company": "Capital Group",
    },
]

CLIENTS_ROW = [
    "Schüco",
    "Alutech",
    "Reynaers",
    "Capital Group",
    "Мосинжпроект",
    "ПАО «Лукойл»",
    "Фонд «Сколково»",
    "Политехстрой",
    "ЖК «Зиларт»",
    "The Ritz-Carlton, Москва",
]

img_counter = 0


def png_for_seed(seed: int) -> Path:
    n = (seed % 53) + 1
    webp = IMG_DIR / f"pptx-{n:02d}.webp"
    png = IMG_DIR / f"pptx-{n:02d}.png"
    return (webp if webp.is_file() else png).resolve()


def next_photo_path() -> Path:
    global img_counter
    p = png_for_seed(img_counter)
    img_counter += 1
    return p


def image_bytes_for_ppt(path: Path) -> bytes:
    """PowerPoint надёжнее принимает PNG; WebP конвертируем при наличии Pillow."""
    data = path.read_bytes()
    if path.suffix.lower() != ".webp":
        return data
    try:
        from PIL import Image

        im = Image.open(BytesIO(data)).convert("RGB")
        buf = BytesIO()
        im.save(buf, format="PNG", optimize=True)
        return buf.getvalue()
    except Exception:
        return data


def replace_picture_blob(shape, path: Path) -> None:
    data = image_bytes_for_ppt(path)
    blip = shape._element.blipFill.blip
    r_id = blip.get(qn("r:embed"))
    part = shape.part.related_part(r_id)
    part._blob = data  # noqa: SLF001


def set_calibri(shape) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            set_calibri(s)
        return
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = "Calibri"


def walk_shapes(slide):
    for sh in slide.shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in sh.shapes:
                yield sub


def patch_slide2(slide) -> None:
    for sh in walk_shapes(slide):
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if "Инжиниринговая компания полного цикла" in t:
            sh.text_frame.text = ABOUT_BODY
        elif t.strip() == "15+":
            sh.text_frame.text = "17+"
        elif t.strip() == "70+":
            sh.text_frame.text = "226+"
        elif t.strip() == "монтажников":
            sh.text_frame.text = "проектов"


def patch_slide3(slide) -> None:
    """Пары заголовок/описание — индексы фигур как в исходном макете."""
    ss = list(slide.shapes)
    titles = [8, 14, 20, 26, 32, 38]
    descs = [10, 16, 22, 28, 34, 40]
    for ti, di, pair in zip(titles, descs, SPECS_LINES, strict=True):
        if ti < len(ss) and ss[ti].has_text_frame:
            ss[ti].text_frame.text = pair[0]
        if di < len(ss) and ss[di].has_text_frame:
            ss[di].text_frame.text = pair[1]


def patch_slide4(slide) -> None:
    for sh in walk_shapes(slide):
        if not sh.has_text_frame:
            continue
        if sh.text_frame.text.strip() == "15+":
            sh.text_frame.text = "17+"


def patch_project_slide(slide, card: dict) -> None:
    ss = list(slide.shapes)
    for idx, val in [(3, card["year"]), (4, card["title"]), (5, card["body"]), (9, card["spec"])]:
        if idx < len(ss) and ss[idx].has_text_frame:
            ss[idx].text_frame.text = val


def patch_slide12(slide) -> None:
    ss = list(slide.shapes)
    slots = [
        (4, 7, 8, 9),
        (11, 14, 15, 16),
        (18, 21, 22, 23),
        (25, 28, 29, 30),
    ]
    for r, tup in zip(REVIEWS, slots, strict=True):
        q_i, a_i, n_i, c_i = tup
        if q_i < len(ss) and ss[q_i].has_text_frame:
            ss[q_i].text_frame.text = f"«{r['quote']}»"
        if a_i < len(ss) and ss[a_i].has_text_frame:
            ss[a_i].text_frame.text = r["initials"]
        if n_i < len(ss) and ss[n_i].has_text_frame:
            ss[n_i].text_frame.text = r["name"]
        if c_i < len(ss) and ss[c_i].has_text_frame:
            ss[c_i].text_frame.text = r["company"]


def patch_slide13(slide) -> None:
    """Плейсхолдеры брендов — фиксированные индексы макета (как в исходном файле)."""
    ss = list(slide.shapes)
    brand_idx = [9, 12, 15, 18, 21, 24, 27, 30, 33, 36]
    for idx, name in zip(brand_idx, CLIENTS_ROW, strict=True):
        if idx < len(ss) and ss[idx].has_text_frame:
            ss[idx].text_frame.text = name
    for sh in walk_shapes(slide):
        if not sh.has_text_frame:
            continue
        if "Ведущие девелоперы" in sh.text_frame.text:
            sh.text_frame.text = "Заказчики и партнёры по проектам и оборудованию"
    if len(ss) > 42 and ss[40].has_text_frame and ss[42].has_text_frame:
        r0 = REVIEWS[0]
        ss[40].text_frame.text = f"«{r0['quote']}»"
        ss[42].text_frame.text = f"{r0['name']}, {r0['company']}"


def patch_slide14(slide) -> None:
    for sh in walk_shapes(slide):
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "888" in t and "(" in t:
            sh.text_frame.text = PHONE
        if t.startswith("info@") or ("@" in t and "a13" in t.lower()):
            sh.text_frame.text = EMAIL
        if t.startswith("Офис:"):
            sh.text_frame.text = f"Офис: {ADDRESS_OFFICE}\nПроизводство: {ADDRESS_PROD}"


def replace_slide_pictures(slide, hero_path: Path | None) -> None:
    """Большое фото — hero_path; остальные PNG — из очереди (тот же стиль, без сдвига рамок)."""
    for sh in walk_shapes(slide):
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        w, h = sh.width, sh.height
        if w >= 9_000_000 and h >= 9_000_000 and hero_path is not None:
            replace_picture_blob(sh, hero_path)
        else:
            replace_picture_blob(sh, next_photo_path())


def main() -> None:
    global img_counter
    if not SRC_PPTX.is_file():
        raise SystemExit(f"Нет файла: {SRC_PPTX}")
    if not IMG_DIR.is_dir():
        raise SystemExit(f"Нет папки изображений: {IMG_DIR}")

    shutil.copy2(SRC_PPTX, OUT_PPTX)
    prs = Presentation(str(OUT_PPTX))
    img_counter = 0

    for si, slide in enumerate(prs.slides):
        if si == 1:
            patch_slide2(slide)
        elif si == 2:
            patch_slide3(slide)
        elif si == 3:
            patch_slide4(slide)
        elif 5 <= si <= 10:
            patch_project_slide(slide, PROJECT_CARDS[si - 5])
        elif si == 11:
            patch_slide12(slide)
        elif si == 12:
            patch_slide13(slide)
        elif si == 13:
            patch_slide14(slide)

        if 5 <= si <= 10:
            hero = png_for_seed(PROJECT_CARDS[si - 5]["img_seed"])
            replace_slide_pictures(slide, hero_path=hero)
        else:
            replace_slide_pictures(slide, hero_path=None)

        for sh in walk_shapes(slide):
            set_calibri(sh)

    prs.save(str(OUT_PPTX))
    print("Saved:", OUT_PPTX)


if __name__ == "__main__":
    main()
