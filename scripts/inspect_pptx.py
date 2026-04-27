# -*- coding: utf-8 -*-
"""Dump PPTX structure for editing."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding="utf-8")

path = Path(r"c:\Users\vladi\Downloads\Untitled.pptx")
out = Path(__file__).resolve().parent / "pptx_dump.txt"
lines: list[str] = []
prs = Presentation(str(path))
for i, slide in enumerate(prs.slides):
    lines.append(f"=== SLIDE {i + 1} ===")
    for j, sh in enumerate(slide.shapes):

        def dump_shape(prefix: str, shape, idx: str) -> None:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    lines.append(f"{prefix}{idx} TEXT: {t[:500]}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                lines.append(f"{prefix}{idx} PICTURE emus {shape.width}x{shape.height}")

        dump_shape("", sh, str(j))
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for k, sub in enumerate(sh.shapes):
                dump_shape("  g", sub, f"{j}.{k}")

out.write_text("\n".join(lines), encoding="utf-8")
print("Wrote", out)
